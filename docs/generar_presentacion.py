"""
Genera la presentación PowerPoint y el guión Word para la defensa del TFM.
Ejecutar con: hm_venv/bin/python docs/generar_presentacion.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pathlib import Path

OUT_DIR = Path(__file__).parent

# ── Paleta de colores ─────────────────────────────────────────────────────────
AZUL_OSCURO  = RGBColor(0x1A, 0x37, 0x5E)   # fondo título
AZUL_MEDIO   = RGBColor(0x1F, 0x4F, 0x8C)   # fondo slides normales
AZUL_CLARO   = RGBColor(0x2E, 0x75, 0xB6)   # acento / bullets
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO   = RGBColor(0xF0, 0xF4, 0xF8)
AMARILLO     = RGBColor(0xFF, 0xC0, 0x00)   # highlight datos importantes

# ── Contenido de cada diapositiva ─────────────────────────────────────────────

SLIDES = [

    # 0 — PORTADA
    {
        "tipo": "portada",
        "titulo": "Sistema RAG y Generador Automático\nde Proyectos Técnicos\npara Reformas de Vehículos",
        "subtitulo": "Trabajo Fin de Máster · Máster en Inteligencia Artificial Aplicada",
        "pie": "Autor: [Tu nombre] · Tutor: [Tutor] · Abril 2026",
    },

    # 1 — PROBLEMA ENCONTRADO
    {
        "tipo": "seccion",
        "num": "01",
        "titulo": "El Problema",
    },
    {
        "tipo": "contenido",
        "titulo": "La tramitación de reformas de vehículos en España",
        "bullets": [
            "El Manual de Reformas de la DGT clasifica 76 tipos de reforma (Fichas CR)",
            "La Vía A exige un Proyecto Técnico completo firmado por ingeniero colegiado",
            "Un proyecto tipo: 30–60 páginas · 4–8 horas de redacción manual",
            "Trabajo repetitivo: compilar normativa, identificar ARs, estructurar secciones",
            "Alta variabilidad en calidad entre ingenieros y talleres",
        ],
        "nota": "💡 El problema no es de volumen solo — es de valor añadido. El ingeniero dedica horas a trabajo mecánico que no aporta criterio técnico diferencial.",
    },
    {
        "tipo": "esquema",
        "titulo": "Flujo actual vs. flujo con el sistema",
        "contenido_izq": [
            "📋 FLUJO ACTUAL",
            "",
            "Ingeniero recibe reforma",
            "↓",
            "Busca manualmente en PDF (76 fichas)",
            "↓",
            "Identifica CRs aplicables",
            "↓",
            "Filtra ARs por categoría de vehículo",
            "↓",
            "Redacta secciones desde plantilla propia",
            "↓",
            "Genera Word manualmente",
            "",
            "⏱ 4–8 horas",
        ],
        "contenido_der": [
            "🤖 CON EL SISTEMA",
            "",
            "Ingeniero introduce datos",
            "↓",
            "RAG recupera fichas CR automáticamente",
            "↓",
            "LLM identifica CRs + filtra ARs",
            "↓",
            "Agentes redactan borrador en paralelo",
            "↓",
            "Ingeniero revisa y aprueba secciones",
            "↓",
            "Sistema genera .docx estructurado",
            "",
            "⏱ < 30 minutos",
        ],
    },

    # 2 — PLANTEAMIENTO DE SOLUCIÓN
    {
        "tipo": "seccion",
        "num": "02",
        "titulo": "Planteamiento de Solución",
    },
    {
        "tipo": "contenido",
        "titulo": "Dos módulos complementarios",
        "bullets": [
            "MÓDULO 1 — Chatbot RAG: consulta del Manual de Reformas en lenguaje natural",
            "MÓDULO 2 — Generador de Proyectos Técnicos Vía A con revisión humana integrada",
            "Base compartida: ChromaDB con 76 fichas CR + preámbulo + Reglamento (UE) 2018/858",
            "Modelos OpenAI: gpt-4o (razonamiento normativo) · gpt-4o-mini (redacción técnica)",
            "Human-in-the-loop: el ingeniero revisa y aprueba cada sección antes del .docx final",
        ],
        "nota": "⚠️ El sistema no reemplaza al ingeniero — elimina el trabajo mecánico y libera tiempo para el criterio técnico y la responsabilidad profesional.",
    },

    # 3 — DOCUMENTOS E INPUTS
    {
        "tipo": "seccion",
        "num": "03",
        "titulo": "Documentos e Inputs",
    },
    {
        "tipo": "contenido",
        "titulo": "Fuentes de conocimiento indexadas",
        "bullets": [
            "Manual de Reformas DGT — Sección I: 76 fichas CR (PDF oficial, ~200 páginas)",
            "Preámbulo del Manual: interpretación de los Actos Reglamentarios (ARs)",
            "Reglamento (UE) 2018/858: definición de categorías de vehículos (M1, M2, N1…)",
            "Keywords de taller: CSV editable con términos coloquiales mapeados a cada CR",
        ],
        "nota": "📂 Inputs del usuario: datos del vehículo (bastidor, categoría, matrícula), datos del ingeniero y del taller, descripción libre de la reforma, CRs opcionales indicados manualmente.",
    },

    # 4 — PARSER
    {
        "tipo": "seccion",
        "num": "04",
        "titulo": "Parser",
    },
    {
        "tipo": "contenido",
        "titulo": "Extracción del PDF con pdfplumber",
        "bullets": [
            "pdfplumber: extracción híbrida — tablas con extract_tables(), texto con extract_text()",
            "Tablas: campo de aplicación por categorías · Actos Reglamentarios con nivel de exigencia",
            "Texto libre: denominación, vía de tramitación, documentación exigible, info adicional",
            "Resultado: 76 objetos JSON con todos los campos estructurados de cada ficha CR",
            "Detección de referencias cruzadas: \"ver también CR X.X\" en informacion_adicional",
        ],
        "nota": "🔧 El reto principal: los PDFs de la DGT no tienen estructura semántica uniforme. Cada ficha tiene variaciones en el formato de las tablas de ARs.",
        "placeholder_imagen": "[INSERTAR: captura del PDF original vs. JSON parseado]",
    },

    # 5 — ENRIQUECIMIENTO
    {
        "tipo": "seccion",
        "num": "05",
        "titulo": "Enriquecimiento",
    },
    {
        "tipo": "contenido",
        "titulo": "Keywords de taller para mejorar el retrieval",
        "bullets": [
            "Problema: el ingeniero busca «turbo» pero la ficha dice «sobrealimentación»",
            "Solución: CSV editable con keywords coloquiales mapeadas a cada CR",
            "Script enriquecimiento.py: merge sin duplicados, validación de columnas, stats",
            "El campo keywords_reformas se inyecta en el chunk de embedding de cada ficha",
            "Actualizable sin re-indexar todo: solo los chunks modificados",
        ],
        "codigo": 'cr,keyword\n2.1,turbocompresor\n2.1,turbo\n8.20,alerón\n8.20,spoiler',
    },

    # 6 — INDEXADO Y BASE VECTORIAL
    {
        "tipo": "seccion",
        "num": "06",
        "titulo": "Indexado y Base Vectorial",
    },
    {
        "tipo": "contenido",
        "titulo": "ChromaDB: tres colecciones con metadatos",
        "bullets": [
            "text-embedding-3-small (OpenAI): vectores de 1.536 dimensiones por chunk",
            "Colección fichas_cr: 76 documentos · metadatos: cr, via_tramitacion, pagina_inicio",
            "Colección preambulo: chunks del preámbulo por apartado temático",
            "Colección reglamento_ue: 8 fragmentos del Reglamento (UE) 2018/858",
            "Retrieval condicional: preámbulo (si query menciona documentación) · reglamento (si menciona categorías)",
            "Filtro por metadatos antes de búsqueda semántica: via_tramitacion = 'A'",
        ],
        "placeholder_imagen": "[INSERTAR: diagrama de las 3 colecciones y flujo de retrieval]",
    },

    # 7 — DOS SERVICIOS
    {
        "tipo": "seccion",
        "num": "07",
        "titulo": "Los Dos Servicios",
    },
    {
        "tipo": "dos_columnas",
        "titulo": "Chatbot RAG · Generador de Proyectos",
        "col_izq_titulo": "💬 CHATBOT RAG",
        "col_izq": [
            "Consulta en lenguaje natural",
            "Historial de conversación (4 turnos)",
            "Filtros: categoría de vehículo · vía",
            "Fuentes citadas en cada respuesta",
            "Retrieval en 3 colecciones",
            "Modelo: gpt-4o-mini",
            "",
            "Ejemplo:",
            "«¿Qué necesito para instalar",
            "un turbo en un M1 Vía A?»",
        ],
        "col_der_titulo": "📄 GENERADOR VÍA A",
        "col_der": [
            "Pipeline multi-agente (LangGraph)",
            "Identificación automática de CRs",
            "Redacción de 10 secciones en paralelo",
            "Human-in-the-loop: revisión sección a sección",
            "Regeneración parcial (solo lo rechazado)",
            "Exportación .docx estructurado",
            "",
            "Modelos:",
            "gpt-4o (Identificador CR)",
            "gpt-4o-mini (Redactores)",
        ],
    },

    # 8 — CARACTERÍSTICAS TÉCNICAS RAG
    {
        "tipo": "seccion",
        "num": "08",
        "titulo": "Características Técnicas del RAG",
    },
    {
        "tipo": "contenido",
        "titulo": "Pipeline RAG con retrieval condicional",
        "bullets": [
            "Pregunta del usuario → embedding → búsqueda coseno en ChromaDB",
            "Retrieval condicional: keywords activan colecciones adicionales (preámbulo, reglamento)",
            "Enriquecimiento de query corta: se concatena con el último turno del historial",
            "Filtro de metadatos previo: reduce el espacio de búsqueda antes del ranking vectorial",
            "Contexto estructurado por secciones: === FICHAS === · === PREÁMBULO === · === REGLAMENTO ===",
            "Temperatura 0 + fuentes citadas → respuestas reproducibles y trazables",
        ],
        "placeholder_imagen": "[INSERTAR: captura del chatbot en funcionamiento]",
    },

    # 9 — ARQUITECTURA
    {
        "tipo": "seccion",
        "num": "09",
        "titulo": "Arquitectura de la Aplicación",
    },
    {
        "tipo": "esquema_arq",
        "titulo": "Frontend · Backend · Base vectorial",
        "lineas": [
            "┌─────────────────────────────────────────────────────────────┐",
            "│   FRONTEND (Streamlit · puerto 8501)                        │",
            "│   Menú central → Chatbot RAG  /  Generador Proyectos        │",
            "└────────────────────────┬────────────────────────────────────┘",
            "                         │ HTTP / REST",
            "┌────────────────────────▼────────────────────────────────────┐",
            "│   BACKEND (FastAPI · puerto 8000)                           │",
            "│   POST /consulta  ·  POST /proyecto-tecnico/iniciar         │",
            "│   POST /proyecto-tecnico/revisar  ·  GET /health            │",
            "│                                                             │",
            "│   RAG chain          LangGraph Grafo                        │",
            "│   retriever.py       identificador_cr → redactores →        │",
            "│                      revision_humana → ensamblador          │",
            "└────────────────────────┬────────────────────────────────────┘",
            "                         │",
            "┌────────────────────────▼────────────────────────────────────┐",
            "│   ChromaDB (local)     OpenAI API          Node.js (.docx)  │",
            "│   fichas_cr            text-embedding-3-small               │",
            "│   preambulo            gpt-4o / gpt-4o-mini                 │",
            "│   reglamento_ue                                             │",
            "└─────────────────────────────────────────────────────────────┘",
        ],
    },
    {
        "tipo": "contenido",
        "titulo": "Endpoints FastAPI y validación Pydantic",
        "bullets": [
            "GET  /health — estado del servicio",
            "POST /consulta — RAG: pregunta + historial + filtros → respuesta + fuentes",
            "GET  /categorias · GET /vias — datos de referencia",
            "POST /proyecto-tecnico/iniciar — lanza el grafo LangGraph",
            "POST /proyecto-tecnico/revisar — reanuda el grafo tras revisión humana",
            "GET  /proyecto-tecnico/descargar/{id} — devuelve el .docx generado",
            "Validación de entrada con Pydantic v2: descripción mínima 20 chars, campos obligatorios",
        ],
    },

    # 10 — FUNCIONAMIENTO FINAL
    {
        "tipo": "seccion",
        "num": "10",
        "titulo": "Funcionamiento Final",
    },
    {
        "tipo": "contenido",
        "titulo": "Pipeline completo del Generador (Vía A)",
        "bullets": [
            "1. Ingeniero introduce datos: vehículo · reforma · taller · componentes",
            "2. Agente 1 (gpt-4o): RAG → identifica CRs · filtra ARs por categoría",
            "3. Agentes 2-4 en paralelo (gpt-4o-mini): redactan 10 secciones",
            "4. Revisión humana: ingeniero aprueba o pide reescritura con motivo",
            "5. Solo se regeneran las secciones rechazadas (no todo el borrador)",
            "6. Agente 5 (Node.js): monta el .docx con portada, índice, tablas y bloques de firma",
        ],
        "placeholder_imagen": "[INSERTAR: captura del paso 3 — revisión de secciones en Streamlit]",
    },
    {
        "tipo": "contenido",
        "titulo": "Tests y calidad del código",
        "bullets": [
            "122 tests unitarios en pytest — 0 llamadas reales a OpenAI ni ChromaDB",
            "test_models.py · test_rag_chain.py · test_rag_retriever.py",
            "test_identificador_cr.py · test_ensamblador.py · test_api.py · test_enriquecimiento.py",
            "FastAPI TestClient + unittest.mock para endpoints y cadena RAG",
            "Despliegue reproducible con Docker: docker-compose up (backend + frontend)",
        ],
    },

    # 11 — TRABAJOS FUTUROS
    {
        "tipo": "seccion",
        "num": "11",
        "titulo": "Trabajos Futuros",
    },
    {
        "tipo": "contenido",
        "titulo": "Líneas de mejora y extensión",
        "bullets": [
            "Vías B, C y D: ampliar el generador más allá de la Vía A",
            "Secciones II-IV del Manual: cálculos justificativos, planos, presupuesto",
            "RAG con reranking: cross-encoder para mejorar la precisión del retrieval",
            "Evaluación automática: RAGAS (faithfulness, answer relevancy, context precision)",
            "Modelos locales (Ollama / Mistral): eliminar dependencia de la API de OpenAI",
            "Integración con firma digital: e-firma del ingeniero directamente en el .docx",
            "Autenticación y multi-tenant: cada despacho con su propia base de proyectos",
        ],
    },

    # 12 — CIERRE
    {
        "tipo": "cierre",
        "titulo": "Gracias",
        "subtitulo": "Sistema RAG y Generador de Proyectos Técnicos para Reformas de Vehículos",
        "repositorio": "github.com/msp40industry-dev/Trabajo_Fin_de_Master",
        "pie": "Turno de preguntas",
    },
]


# ── Helpers de formato PowerPoint ─────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=BLANCO,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=17, color=BLANCO, marker_color=AMARILLO):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        # Marcador de bullet
        run_mark = p.add_run()
        run_mark.text = "▸  "
        run_mark.font.size = Pt(font_size)
        run_mark.font.color.rgb = marker_color
        run_mark.font.bold = True
        # Texto
        run_text = p.add_run()
        run_text.text = bullet
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
    return txBox


def add_rect(slide, color: RGBColor, left, top, width, height):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Generador de slides ───────────────────────────────────────────────────────

def hacer_portada(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, AZUL_OSCURO)
    W = prs.slide_width
    H = prs.slide_height

    # Franja lateral izquierda
    add_rect(slide, AZUL_CLARO, 0, 0, Inches(0.35), H)

    # Línea decorativa horizontal
    add_rect(slide, AMARILLO, Inches(0.5), Inches(3.8), W - Inches(1), Inches(0.06))

    # Título
    add_text_box(slide, data["titulo"],
                 Inches(0.6), Inches(1.2), W - Inches(1.2), Inches(2.5),
                 font_size=30, bold=True, color=BLANCO)

    # Subtítulo
    add_text_box(slide, data["subtitulo"],
                 Inches(0.6), Inches(4.0), W - Inches(1.2), Inches(0.7),
                 font_size=16, color=RGBColor(0xB0, 0xC8, 0xE8))

    # Pie
    add_text_box(slide, data["pie"],
                 Inches(0.6), H - Inches(1.0), W - Inches(1.2), Inches(0.5),
                 font_size=13, color=RGBColor(0x90, 0xA8, 0xC0))


def hacer_seccion(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_OSCURO)
    W = prs.slide_width
    H = prs.slide_height

    # Número grande decorativo
    add_text_box(slide, data["num"],
                 Inches(0.3), Inches(0.5), Inches(2.5), Inches(3),
                 font_size=120, bold=True,
                 color=RGBColor(0x2E, 0x55, 0x8C))  # azul semitransparente

    # Línea
    add_rect(slide, AMARILLO, Inches(0.5), Inches(3.5), Inches(5), Inches(0.07))

    # Título
    add_text_box(slide, data["titulo"],
                 Inches(0.5), Inches(3.7), W - Inches(1), Inches(1.5),
                 font_size=38, bold=True, color=BLANCO)


def hacer_contenido(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_MEDIO)
    W = prs.slide_width
    H = prs.slide_height

    # Franja superior
    add_rect(slide, AZUL_OSCURO, 0, 0, W, Inches(1.15))

    # Título
    add_text_box(slide, data["titulo"],
                 Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.9),
                 font_size=22, bold=True, color=BLANCO)

    # Bullets
    bullets = data.get("bullets", [])
    if bullets:
        add_bullet_box(slide, bullets,
                       Inches(0.4), Inches(1.3), W - Inches(0.8), H - Inches(2.5))

    # Nota al pie
    nota = data.get("nota")
    if nota:
        add_rect(slide, AZUL_OSCURO, Inches(0.3), H - Inches(1.3), W - Inches(0.6), Inches(1.0))
        add_text_box(slide, nota,
                     Inches(0.45), H - Inches(1.25), W - Inches(0.9), Inches(0.9),
                     font_size=13, color=RGBColor(0xB0, 0xD8, 0xFF))

    # Placeholder imagen
    placeholder = data.get("placeholder_imagen")
    if placeholder:
        add_rect(slide, RGBColor(0x15, 0x35, 0x60),
                 W - Inches(4.2), Inches(1.3), Inches(3.9), Inches(3.5))
        add_text_box(slide, placeholder,
                     W - Inches(4.1), Inches(2.7), Inches(3.7), Inches(0.7),
                     font_size=11, color=RGBColor(0x80, 0xA8, 0xCC), align=PP_ALIGN.CENTER)

    # Bloque de código
    codigo = data.get("codigo")
    if codigo:
        add_rect(slide, RGBColor(0x0F, 0x20, 0x38),
                 Inches(0.4), H - Inches(2.2), W - Inches(0.8), Inches(1.4))
        add_text_box(slide, codigo,
                     Inches(0.55), H - Inches(2.1), W - Inches(1.0), Inches(1.2),
                     font_size=14, color=RGBColor(0x80, 0xFF, 0x80))


def hacer_esquema(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_MEDIO)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, AZUL_OSCURO, 0, 0, W, Inches(1.15))
    add_text_box(slide, data["titulo"],
                 Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.9),
                 font_size=22, bold=True, color=BLANCO)

    mid = W // 2

    # Columna izquierda
    add_rect(slide, AZUL_OSCURO, Inches(0.3), Inches(1.25), mid - Inches(0.5), H - Inches(1.5))
    izq_text = "\n".join(data.get("contenido_izq", []))
    add_text_box(slide, izq_text,
                 Inches(0.45), Inches(1.35), mid - Inches(0.7), H - Inches(1.7),
                 font_size=14, color=BLANCO)

    # Columna derecha
    add_rect(slide, RGBColor(0x10, 0x30, 0x58), mid + Inches(0.2), Inches(1.25),
             mid - Inches(0.5), H - Inches(1.5))
    der_text = "\n".join(data.get("contenido_der", []))
    add_text_box(slide, der_text,
                 mid + Inches(0.35), Inches(1.35), mid - Inches(0.7), H - Inches(1.7),
                 font_size=14, color=BLANCO)


def hacer_dos_columnas(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_MEDIO)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, AZUL_OSCURO, 0, 0, W, Inches(1.15))
    add_text_box(slide, data["titulo"],
                 Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.9),
                 font_size=22, bold=True, color=BLANCO)

    mid = W // 2

    # Izquierda
    add_rect(slide, AZUL_OSCURO, Inches(0.3), Inches(1.25), mid - Inches(0.5), H - Inches(1.5))
    add_text_box(slide, data["col_izq_titulo"],
                 Inches(0.45), Inches(1.35), mid - Inches(0.7), Inches(0.5),
                 font_size=15, bold=True, color=AMARILLO)
    izq_text = "\n".join(data.get("col_izq", []))
    add_text_box(slide, izq_text,
                 Inches(0.45), Inches(1.95), mid - Inches(0.7), H - Inches(2.2),
                 font_size=13, color=BLANCO)

    # Derecha
    add_rect(slide, RGBColor(0x10, 0x30, 0x58), mid + Inches(0.2), Inches(1.25),
             mid - Inches(0.5), H - Inches(1.5))
    add_text_box(slide, data["col_der_titulo"],
                 mid + Inches(0.35), Inches(1.35), mid - Inches(0.7), Inches(0.5),
                 font_size=15, bold=True, color=AMARILLO)
    der_text = "\n".join(data.get("col_der", []))
    add_text_box(slide, der_text,
                 mid + Inches(0.35), Inches(1.95), mid - Inches(0.7), H - Inches(2.2),
                 font_size=13, color=BLANCO)


def hacer_esquema_arq(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_OSCURO)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, AZUL_CLARO, 0, 0, W, Inches(1.0))
    add_text_box(slide, data["titulo"],
                 Inches(0.4), Inches(0.15), W - Inches(0.8), Inches(0.75),
                 font_size=22, bold=True, color=BLANCO)

    texto = "\n".join(data["lineas"])
    add_text_box(slide, texto,
                 Inches(0.3), Inches(1.1), W - Inches(0.6), H - Inches(1.3),
                 font_size=12, color=RGBColor(0xB0, 0xD8, 0xFF))


def hacer_cierre(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_OSCURO)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, AZUL_CLARO, 0, 0, W, Inches(0.35))
    add_rect(slide, AZUL_CLARO, 0, H - Inches(0.35), W, Inches(0.35))

    add_text_box(slide, data["titulo"],
                 Inches(0.5), Inches(1.5), W - Inches(1), Inches(1.5),
                 font_size=60, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    add_rect(slide, AMARILLO, Inches(2), Inches(3.2), W - Inches(4), Inches(0.07))

    add_text_box(slide, data["subtitulo"],
                 Inches(0.5), Inches(3.4), W - Inches(1), Inches(0.8),
                 font_size=16, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

    add_text_box(slide, f"🔗 {data['repositorio']}",
                 Inches(0.5), Inches(4.4), W - Inches(1), Inches(0.5),
                 font_size=14, color=AMARILLO, align=PP_ALIGN.CENTER)

    add_text_box(slide, data["pie"],
                 Inches(0.5), H - Inches(1.2), W - Inches(1), Inches(0.6),
                 font_size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)


# ── Construir PowerPoint ──────────────────────────────────────────────────────

def construir_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    handlers = {
        "portada":     hacer_portada,
        "seccion":     hacer_seccion,
        "contenido":   hacer_contenido,
        "esquema":     hacer_esquema,
        "dos_columnas": hacer_dos_columnas,
        "esquema_arq": hacer_esquema_arq,
        "cierre":      hacer_cierre,
    }

    for slide_data in SLIDES:
        tipo = slide_data["tipo"]
        handlers[tipo](prs, slide_data)

    out = OUT_DIR / "Defensa_TFM.pptx"
    prs.save(str(out))
    print(f"✅ PowerPoint guardado: {out}")
    return out


# ── Construir guión Word ──────────────────────────────────────────────────────

GUION = [
    {
        "slide": "PORTADA",
        "tiempo": "0:00 – 0:45",
        "titulo": "Presentación inicial",
        "texto": (
            "Buenos días. Me llamo [nombre] y voy a presentar mi Trabajo Fin de Máster: "
            "un sistema de consulta y generación automática de proyectos técnicos para la "
            "tramitación de reformas de vehículos en España.\n\n"
            "La idea nace de una necesidad real que observé en el entorno profesional, "
            "y que voy a explicar a continuación."
        ),
    },
    {
        "slide": "DIAPOSITIVA 1 — El Problema",
        "tiempo": "0:45 – 2:30",
        "titulo": "Contexto normativo y problema",
        "texto": (
            "En España, cualquier modificación de un vehículo que afecte a sus características "
            "técnicas debe tramitarse según el Manual de Reformas de la DGT. Este manual clasifica "
            "76 tipos distintos de reforma en lo que se llaman Fichas CR.\n\n"
            "La vía más exigente —la Vía A— requiere un Proyecto Técnico completo: un documento "
            "de 30 a 60 páginas que debe firmar un ingeniero colegiado. Elaborarlo a mano "
            "lleva entre 4 y 8 horas.\n\n"
            "Y el problema no es solo el tiempo: la mayor parte de ese trabajo es compilación "
            "mecánica de normativa — identificar qué CRs aplican, qué Actos Reglamentarios "
            "hay que cumplir según la categoría del vehículo, estructurar el documento... "
            "Trabajo repetitivo que no aporta criterio técnico diferencial.\n\n"
            "[→ MOSTRAR ESQUEMA: flujo actual vs. flujo con el sistema]"
        ),
    },
    {
        "slide": "DIAPOSITIVA 2 — Planteamiento de solución",
        "tiempo": "2:30 – 4:00",
        "titulo": "La solución: dos módulos complementarios",
        "texto": (
            "La solución que propongo tiene dos módulos complementarios.\n\n"
            "El primero es un chatbot de consulta basado en RAG — Retrieval-Augmented Generation. "
            "Permite al ingeniero o al taller hacer preguntas en lenguaje natural sobre el Manual "
            "de Reformas, sin tener que navegar por el PDF.\n\n"
            "El segundo es el generador automático de proyectos técnicos Vía A. A partir de los "
            "datos del vehículo y la descripción de la reforma, el sistema identifica los CRs, "
            "redacta el borrador por secciones y permite que el ingeniero lo revise y apruebe "
            "antes de exportarlo a Word.\n\n"
            "Importante: el sistema no elimina la responsabilidad del ingeniero. Elimina el "
            "trabajo mecánico y le deja tiempo para el criterio técnico."
        ),
    },
    {
        "slide": "DIAPOSITIVA 3 — Documentos e Inputs",
        "tiempo": "4:00 – 5:15",
        "titulo": "Fuentes de conocimiento y datos de entrada",
        "texto": (
            "El sistema se alimenta de tres fuentes de conocimiento. La principal es la Sección I "
            "del Manual de Reformas —las 76 fichas CR—. La segunda es el preámbulo del manual, "
            "que explica cómo interpretar los niveles de exigencia de los Actos Reglamentarios. "
            "Y la tercera es el Reglamento (UE) 2018/858, que define las categorías de vehículos.\n\n"
            "Además hay un cuarto input configurable: un CSV de keywords de taller. Esto resuelve "
            "un problema de vocabulario — el ingeniero busca 'turbo' pero la ficha dice "
            "'sobrealimentación'. El CSV mapea términos coloquiales a cada CR.\n\n"
            "Por parte del usuario, el sistema recibe los datos del vehículo, del ingeniero, "
            "del taller y la descripción libre de la reforma."
        ),
    },
    {
        "slide": "DIAPOSITIVA 4 — Parser",
        "tiempo": "5:15 – 6:30",
        "titulo": "Extracción estructurada del PDF",
        "texto": (
            "El pipeline empieza con la extracción del PDF del Manual de Reformas usando pdfplumber.\n\n"
            "La estrategia es híbrida: las tablas de campo de aplicación y de Actos Reglamentarios "
            "se extraen con extract_tables(), porque tienen estructura tabular. Los campos de texto "
            "libre —descripción, documentación exigible, información adicional— se extraen con "
            "extract_text().\n\n"
            "El resultado son 76 objetos JSON totalmente estructurados, con todos los campos de "
            "cada ficha CR. El mayor reto técnico aquí fue que los PDFs de la DGT no tienen "
            "formato uniforme entre fichas — cada una tiene pequeñas variaciones en la disposición "
            "de las tablas de ARs."
        ),
    },
    {
        "slide": "DIAPOSITIVA 5 — Enriquecimiento",
        "tiempo": "6:30 – 7:30",
        "titulo": "Keywords de taller para mejorar el retrieval",
        "texto": (
            "Una vez parseadas las fichas, hay un paso de enriquecimiento opcional. "
            "Se parte de un CSV editable con pares CR-keyword. El script enriquecimiento.py "
            "hace el merge, deduplica en case-insensitive y actualiza el campo keywords_reformas "
            "de cada ficha.\n\n"
            "Esto mejora el retrieval semántico porque el chunk de embedding incluye tanto el "
            "vocabulario normativo como el vocabulario de taller. La actualización no requiere "
            "re-indexar todo — solo los chunks modificados."
        ),
    },
    {
        "slide": "DIAPOSITIVA 6 — Indexado y Base Vectorial",
        "tiempo": "7:30 – 9:00",
        "titulo": "ChromaDB: tres colecciones con retrieval condicional",
        "texto": (
            "El indexado usa el modelo text-embedding-3-small de OpenAI, que genera vectores "
            "de 1.536 dimensiones. Tenemos tres colecciones en ChromaDB: fichas_cr, preambulo "
            "y reglamento_ue.\n\n"
            "El retrieval no es siempre igual: hay condiciones. Si la query menciona 'documentación', "
            "'proyecto técnico' o similares, se activa el retrieval en el preámbulo. Si menciona "
            "categorías de vehículo como M1 o N1, se activa el retrieval en el reglamento. "
            "Y antes de la búsqueda vectorial se pueden aplicar filtros por metadatos — por vía "
            "de tramitación, por ejemplo — para reducir el espacio de búsqueda.\n\n"
            "[→ MOSTRAR DIAGRAMA de las tres colecciones]"
        ),
    },
    {
        "slide": "DIAPOSITIVA 7 — Los Dos Servicios",
        "tiempo": "9:00 – 10:30",
        "titulo": "Chatbot RAG vs. Generador de Proyectos",
        "texto": (
            "Veamos los dos servicios con más detalle.\n\n"
            "El chatbot RAG es el más sencillo: el usuario hace una pregunta, el retriever "
            "recupera los chunks relevantes de ChromaDB, se construye el contexto estructurado "
            "y gpt-4o-mini genera la respuesta citando las fuentes. Mantiene historial de hasta "
            "4 turnos y permite filtrar por categoría de vehículo o vía de tramitación.\n\n"
            "El generador es más complejo. Es un grafo multi-agente implementado con LangGraph. "
            "El agente Identificador usa gpt-4o —el más potente— para la tarea crítica de "
            "razonamiento normativo. Los tres agentes redactores usan gpt-4o-mini en paralelo "
            "para generar las secciones del documento. Y hay un punto de interrupción —"
            "human-in-the-loop— donde el ingeniero revisa y aprueba.\n\n"
            "[→ MOSTRAR captura del chatbot y del generador en Streamlit]"
        ),
    },
    {
        "slide": "DIAPOSITIVA 8 — Características Técnicas del RAG",
        "tiempo": "10:30 – 12:00",
        "titulo": "RAG con retrieval condicional y contexto estructurado",
        "texto": (
            "El pipeline RAG tiene algunas características que vale la pena destacar.\n\n"
            "Primero, el retrieval condicional: no siempre se consultan las tres colecciones. "
            "Esto reduce ruido y mejora la precisión. Segundo, el enriquecimiento de queries "
            "cortas: si la pregunta tiene menos de 6 palabras, se concatena con el turno anterior "
            "del historial para tener más contexto de búsqueda.\n\n"
            "El contexto que recibe el LLM está estructurado en secciones bien delimitadas — "
            "fichas, preámbulo, reglamento — para que el modelo sepa de dónde viene cada dato.\n\n"
            "Y la temperatura es 0 en todos los agentes: máxima reproducibilidad, sin "
            "variabilidad aleatoria. Para documentos con efectos legales, esto es importante."
        ),
    },
    {
        "slide": "DIAPOSITIVA 9 — Arquitectura",
        "tiempo": "12:00 – 14:00",
        "titulo": "Frontend, Backend, Base Vectorial",
        "texto": (
            "La arquitectura tiene dos capas separadas.\n\n"
            "El frontend es una aplicación Streamlit con un menú central que permite elegir entre "
            "el chatbot y el generador. Corre en el puerto 8501.\n\n"
            "El backend es FastAPI en el puerto 8000. Expone los endpoints del RAG y los del "
            "generador de proyectos. Toda la validación de entrada la hace Pydantic v2 — si falta "
            "un campo obligatorio o la descripción de la reforma es demasiado corta, la API devuelve "
            "422 antes de llegar al LLM.\n\n"
            "Los endpoints del generador son dos: /iniciar lanza el grafo LangGraph, que se "
            "interrumpe antes de la revisión. /revisar recibe las aprobaciones del ingeniero "
            "y reanuda el grafo. Si hay secciones rechazadas, solo se regeneran esas — el resto "
            "se mantiene.\n\n"
            "El despliegue es Docker con docker-compose: dos contenedores, backend y frontend, "
            "con ChromaDB montado como volumen local."
        ),
    },
    {
        "slide": "DIAPOSITIVA 10 — Funcionamiento Final",
        "tiempo": "14:00 – 16:30",
        "titulo": "Pipeline completo y tests",
        "texto": (
            "El pipeline completo del generador tiene 6 pasos.\n\n"
            "El ingeniero introduce los datos del vehículo y la reforma. El agente Identificador "
            "CR —con gpt-4o— consulta ChromaDB, identifica las fichas aplicables y filtra los "
            "ARs por categoría. Los tres redactores —con gpt-4o-mini— generan las secciones en "
            "paralelo. El ingeniero revisa sección por sección y puede aprobar o pedir reescritura "
            "con un motivo concreto. Solo se regeneran las secciones rechazadas. Y finalmente el "
            "ensamblador genera el .docx con portada, índice, tablas normativas y espacio para firma.\n\n"
            "[→ MOSTRAR captura del paso de revisión en Streamlit]\n\n"
            "En cuanto a calidad del código: hay 122 tests unitarios que validan toda la lógica "
            "pura sin llamadas reales a OpenAI ni ChromaDB. Se usa FastAPI TestClient y "
            "unittest.mock para los endpoints."
        ),
    },
    {
        "slide": "DIAPOSITIVA 11 — Trabajos Futuros",
        "tiempo": "16:30 – 18:30",
        "titulo": "Líneas de mejora y extensión",
        "texto": (
            "Hay varias líneas de trabajo futuro identificadas.\n\n"
            "La más natural es extender el generador a las vías B, C y D, y añadir las "
            "secciones más complejas del proyecto —cálculos justificativos, planos, presupuesto—.\n\n"
            "En cuanto a la calidad del retrieval, sería interesante añadir un reranker "
            "cross-encoder y evaluar con métricas de RAGAS: faithfulness, answer relevancy "
            "y context precision.\n\n"
            "A nivel de infraestructura, la alternativa más impactante sería sustituir la API "
            "de OpenAI por modelos locales con Ollama — para entornos donde la privacidad de "
            "los datos del vehículo sea un requisito estricto.\n\n"
            "Y a nivel de producto, la integración con firma digital y un sistema de "
            "autenticación multi-tenant permitirían que varios despachos usaran el sistema "
            "de forma independiente."
        ),
    },
    {
        "slide": "CIERRE",
        "tiempo": "18:30 – 20:00",
        "titulo": "Conclusión y agradecimientos",
        "texto": (
            "Para cerrar: el sistema demuestra que es posible automatizar la parte mecánica de "
            "la redacción de proyectos técnicos de reforma sin eliminar la supervisión "
            "profesional del ingeniero.\n\n"
            "Los dos módulos —chatbot RAG y generador— comparten la misma base de conocimiento "
            "y se despliegan con Docker en cualquier máquina. El tiempo de generación de un "
            "borrador completo está por debajo de 3 minutos, con un coste por proyecto inferior "
            "a 0,50 dólares de API.\n\n"
            "El repositorio está disponible en GitHub. Muchas gracias por su atención. "
            "Quedo a disposición del tribunal para las preguntas."
        ),
    },
    {
        "slide": "TURNO DE PREGUNTAS — Posibles preguntas del tribunal",
        "tiempo": "20:00 – 30:00",
        "titulo": "Preguntas frecuentes y respuestas preparadas",
        "texto": (
            "── ¿Por qué RAG en lugar de fine-tuning?\n"
            "El fine-tuning requiere datos etiquetados, es costoso de reentrenar cuando cambia "
            "la normativa, y no garantiza trazabilidad de fuentes. RAG permite actualizar el "
            "corpus sin tocar el modelo y citar la fuente exacta de cada respuesta. Para un "
            "dominio normativo que se actualiza periódicamente, RAG es la opción correcta.\n\n"
            "── ¿Cómo gestionáis las alucinaciones del LLM?\n"
            "Tres mecanismos: temperatura 0, inyección explícita de los ARs y fichas en el "
            "prompt (el modelo redacta a partir de datos, no de memoria), y revisión humana "
            "obligatoria antes de exportar. Ningún documento sale sin que el ingeniero lo haya "
            "aprobado sección a sección.\n\n"
            "── ¿Qué pasa si el ingeniero rechaza muchas secciones?\n"
            "El grafo LangGraph permite regeneración parcial indefinida. Solo se regeneran las "
            "secciones marcadas para reescritura, con el motivo del ingeniero inyectado en el "
            "prompt. El resto se conserva. No hay límite de iteraciones, aunque cada regeneración "
            "tiene coste de API.\n\n"
            "── ¿Por qué ChromaDB y no Pinecone o Qdrant?\n"
            "Para el volumen actual —76 fichas más unos cientos de chunks adicionales— ChromaDB "
            "en modo embebido es suficiente y no requiere servicio externo. Si el corpus crece "
            "significativamente o se necesita multi-tenant, la migración a Qdrant sería la opción "
            "más natural por su soporte de filtros por metadatos similares.\n\n"
            "── ¿Habéis evaluado la calidad de las respuestas del chatbot?\n"
            "Se ha validado cualitativamente con casos de uso reales. La evaluación con RAGAS "
            "está identificada como trabajo futuro — es precisamente la métrica que usaría para "
            "una evaluación rigurosa: faithfulness, answer relevancy y context precision.\n\n"
            "── ¿Por qué Node.js para el .docx y no python-docx?\n"
            "La librería docx de Node.js tiene soporte más completo de OOXML en el momento del "
            "desarrollo: cabeceras y pies de página con numeración, índices con hipervínculos, "
            "estilos de tabla complejos. python-docx tiene limitaciones en esas áreas que "
            "harían más difícil generar el documento con la maquetación requerida."
        ),
    },
]


def construir_guion():
    doc = Document()

    # Estilos generales
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DPt(11)

    # Portada del guión
    titulo = doc.add_heading("Guión de Defensa — TFM", level=0)
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitulo = doc.add_paragraph(
        "Sistema RAG y Generador Automático de Proyectos Técnicos para Reformas de Vehículos"
    )
    subtitulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitulo.runs[0].bold = True

    doc.add_paragraph(
        "Duración total: 20 minutos de presentación + 10 minutos de preguntas"
    ).alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("─" * 60)

    intro = doc.add_paragraph()
    intro.add_run("INSTRUCCIONES DE USO:\n").bold = True
    intro.add_run(
        "Este guión indica el tiempo aproximado por diapositiva y el texto orientativo "
        "a desarrollar. No es un texto para leer literalmente — es un apoyo para "
        "mantener el ritmo y no olvidar los puntos clave. Las indicaciones entre "
        "[corchetes] son acciones o transiciones."
    )

    doc.add_page_break()

    for bloque in GUION:
        # Encabezado de sección
        h = doc.add_heading(bloque["slide"], level=1)
        h.runs[0].font.color.rgb = DRGBColor(0x1A, 0x37, 0x5E)

        # Tiempo y título
        p_meta = doc.add_paragraph()
        p_meta.add_run("⏱ Tiempo: ").bold = True
        p_meta.add_run(bloque["tiempo"] + "   ")
        p_meta.add_run("Punto clave: ").bold = True
        p_meta.add_run(bloque["titulo"])

        doc.add_paragraph("─" * 40)

        # Texto del guión
        for parrafo in bloque["texto"].split("\n\n"):
            p = doc.add_paragraph(parrafo)
            p.paragraph_format.space_after = DPt(8)

        doc.add_paragraph("")  # espacio entre bloques

    out = OUT_DIR / "Guion_Defensa_TFM.docx"
    doc.save(str(out))
    print(f"✅ Guión Word guardado: {out}")
    return out


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    construir_pptx()
    construir_guion()
    print("\n✅ Ambos ficheros generados en docs/")
